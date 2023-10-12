import streamlit as st
from pypdf import PdfReader 
# import pdfplumber
from pptx import Presentation
import pinecone
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.vectorstores import Qdrant, Pinecone
import qdrant_client
from qdrant_client import QdrantClient, models
from langchain.chains import RetrievalQA
from langchain.chains.summarize import load_summarize_chain
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
# from mv_scraper import scrape
import os
import pickle
import openai
import langchain

langchain.debug = True
# langchain.verbose = False

client = qdrant_client.QdrantClient(
    url="https://87c999b4-f59f-4573-881e-f0834ef04819.us-east4-0.gcp.cloud.qdrant.io:6333", 
    api_key="9Ay5TA6OX1bKPAVEkbg-vYn60-Z4eBholuhqMvhOmvDmgOH4dhtBhw",
)

collection_config = qdrant_client.http.models.VectorParams(
        size=1536, # 768 for instructor-xl, 1536 for OpenAI
        distance=qdrant_client.http.models.Distance.COSINE
    )

client.recreate_collection(
    collection_name="mv",
    vectors_config=collection_config
)

openai.api_key = st.secrets["OPENAI_API_KEY"]

def get_pdf_text(file):
    try:
        print(f"Processing PDF file: {file}")
        pdf_reader = PdfReader(file)
        text = ""
        for page_number in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_number]
            text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error processing PDF file: {e}")
        return ""

def get_ppt_text(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    index_name = 'multiply-ventures'
    embeddings = OpenAIEmbeddings()

    # # Check if the index exists
    # if index_name not in pinecone.list_indexes():
    #     # If not, create the index
    #     pinecone.create_index(index_name, metric='cosine', dimension=1536)

    # # Load the index
    # index = pinecone.Index(index_name=index_name)

    # Initialize the Pinecone vector store with the index
    vectorstore = Qdrant(
    client=client, collection_name="mv", 
    embeddings=embeddings,)
    vectorstore.add_texts(text_chunks)
    # # Generate embeddings and store them in Qdrant
    # for i, text_chunk in enumerate(text_chunks):
    #     vector = embeddings.embed_query(text_chunk)
    #     metadata = {"text": text_chunk}
    #     vectorstore.upsert_points(
    #         collection_name="my_collection",
    #         points=models.Batch(
    #             ids=[str(i)],
    #             vectors=[vector],
    #             payloads=[models.Payload({"text": metadata})],
    #         )
    #     )

    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model="gpt-3.5-turbo",max_tokens=1000,temperature=0.4)
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = RetrievalQA.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory,
        verbose=True
    )
    return conversation_chain

def handle_userinput(user_question, conversation_chain, chat_history):
    response = conversation_chain({'query': user_question})
    chat_history.append({'query': user_question, 'response': response['chat_history'][-1]})
    return chat_history

from langchain.chains.mapreduce import MapReduceChain

def generate_memo(raw_text):    
    topic = "Evaluate the following venture. Think step by step as a Venture Capitalist. Format should be: Name of the company, what they do, what is their product, what is their market, what is their business model, what is their traction, what is their team, what is their competition, what is their financials, what is their valuation, what is your recommendation."
    
    text_chunks = get_text_chunks(raw_text)
    vectorstore = get_vectorstore(text_chunks)
    prompt_template = """
    You are a seasoned venture capitalist evaluating potential investment opportunities. 
    Your task is to critically evaluate the following venture based on the information provided and your expertise. 
    Write a detailed 1000 word memo to your investment committee evaluating the following venture, use the context below to write a memo about the topic below:
    DO NOT USE WORDS THESE WORDS PROVIDED BELOW MEMO.
    YOUR NAME 
    SUBJECT.To: Investment Committee 
    From: [Your Name] 
    Date: [Date]. 
    Context: {context}
    Topic: {topic}
    Memo:"""
    
    PROMPT = PromptTemplate(template=prompt_template, input_variables=["context", "topic"], validate_template=False)

    llm = ChatOpenAI(model="gpt-3.5-turbo-16k", max_tokens=5000, temperature=0.4)
    chain = LLMChain(llm=llm, prompt=PROMPT)
    docs = vectorstore.similarity_search(topic, k=1)
    context = "\n".join([doc.page_content for doc in docs])
    # inputs = [{"context": context, "topic": "Evaluate the following venture."} for doc in docs]
    # Define StuffDocumentsChain
    
    # Apply the chain to the inputs
    inputs = [{"context": context, "topic": "Evaluate the following venture.", "input_documents": docs} for doc in docs]  # Add "input_documents" to the inputs
    stuff_chain = StuffDocumentsChain(
        llm_chain=chain, document_variable_name="context"
    )
    memos = stuff_chain.apply(inputs)
    print(memos)
    output_texts = [memo['output_text'] for memo in memos]
    combined_output_text = "\n".join(output_texts)
    return combined_output_text
    # memo_docs = [Document(page_content=memo['context']) for memo in memos]
    # # Define the map_chain
    # map_template = """The following is a set of documents
    # {docs}
    # Based on this list of docs, please identify the main themes: 
    # Name of the company, what they do, what is their product, what is their market, what is their business model, what is their traction, what is their team, what is their competition, what is their financials, what is their valuation, what is your recommendation.
    # Helpful Answer:"""
    # map_prompt = PromptTemplate.from_template(map_template)
    # map_chain = LLMChain(llm=llm, prompt=map_prompt)

    # # Define the reduce_chain
    # reduce_template = """The following is set of summaries:
    # {doc_summaries}
    # Take these and distill it into a final, consolidated summary of the main themes. 
    # Helpful Answer:"""
    # reduce_prompt = PromptTemplate.from_template(reduce_template)
    # reduce_chain = LLMChain(llm=llm, prompt=reduce_prompt)

    # # Define the combine_documents_chain
    # combine_documents_chain = StuffDocumentsChain(
    #     llm_chain=reduce_chain, document_variable_name="doc_summaries"
    # )

    # # Define the reduce_documents_chain
    # reduce_documents_chain = ReduceDocumentsChain(
    #     combine_documents_chain=combine_documents_chain,
    #     collapse_documents_chain=combine_documents_chain,
    #     token_max=4000,
    # )

    # # Define the map_reduce_chain
    # map_reduce_chain = MapReduceDocumentsChain(
    #     llm_chain=map_chain,
    #     reduce_documents_chain=reduce_documents_chain,
    #     document_variable_name="docs",
    #     return_intermediate_steps=True,
    # )

    # # Split each document into chunks
    # text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
    # split_docs = text_splitter.split_documents(memo_docs)

    # # Run the map_reduce_chain on the chunked documents
    # result = map_reduce_chain.run(split_docs)

    # print(result)
    # return result
    # # Apply the chain to the inputs
    # memos = stuff_chain.apply(inputs)
    
    # # # Print the first memo to check its structure
    # # print(memos[0])
    
    # # # Convert memos to Document objects
    # memo_docs = [Document(page_content=memo['context']) for memo in memos]
    # # Initialize a TextSplitter with a specified chunk size
    # chunk_size = 1250
    # text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=100)# Adjust this value based on your needs

    # # Split each document into chunks
    # chunked_memo_docs = []
    # for doc in memo_docs:
    #     chunks = text_splitter.split_text(doc.page_content)
    #     chunked_memo_docs.extend([Document(page_content=chunk) for chunk in chunks])

    # # Load the refine summarization chain
    # summarize_chain = load_summarize_chain(llm, chain_type="")

    # # Run the summarization chain on the chunked documents
    # result = []
    # for chunk in chunked_memo_docs:
    #     chunk_result = summarize_chain.run([chunk])
    #     result.append(chunk_result)

    # # Combine the results
    # combined_result = ' '.join(result)

    # print(combined_result)
    # return combined_result
    # # # Load the refine summarization chain
    # # summarize_chain = load_summarize_chain(llm, chain_type="stuff")

    # # # Run the summarization chain on the memos
    # # result = summarize_chain.run(memo_docs)

    # # print(result)
    # # return result
    
def app():
    chat_history = []
    
    st.title('Multiple Ventures')

    file = st.file_uploader("Upload a file", type=['pdf', 'pptx'])
    
    # if st.button('Scrape URLs'):
    #     urls = st.text_input("Enter a list of URLs (comma-separated): ")
    #     if urls:
    #         urls = urls.split(",")
    #         urls = [url.strip() for url in urls]
    #         output_file = "output.txt"
    #         st.text('Starting to scrape URLs...')
    #         messages = scrape(urls, output_file=output_file)
    #         for message in messages:
    #             st.write(message)
    #         st.text('Scraping completed.')
    #         st.success('URLs scraped successfully. Data saved to output.txt.')
            
    if file is not None:
        with st.status('Processing file...'):
            if file.type == 'application/pdf':
                raw_text = get_pdf_text(file)
            elif file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                raw_text = get_ppt_text(file)
            else:
                st.error('Invalid file type. Please upload either a "pdf" or "pptx" file.')

            text_chunks = get_text_chunks(raw_text)
            vectorstore = get_vectorstore(text_chunks)
            st.success('Embeddings Generated Successfully.')
            conversation_chain = get_conversation_chain(vectorstore)
            st.success('File uploaded successfully.')

    if st.button('Generate Memo'):    
        memo = generate_memo(raw_text)
        if memo is not None:
            st.write(memo)
    
    user_question = st.chat_input('Enter your question:')
    if user_question:
        if conversation_chain is not None:  # Check if conversation_chain is not None before handling user input
            chat_history = handle_userinput(user_question, conversation_chain, chat_history)
            for interaction in chat_history:
                with st.chat_message('user'):
                    st.write(interaction['query'])
                with st.chat_message('bot'):
                    st.write(interaction['response'].content)
        else:
            st.error('Please upload a file before asking a question.')
        # Add a button to generate a memo


if __name__ == '__main__':
    app()